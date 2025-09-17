import pytest
from unittest.mock import MagicMock, patch

from pptx import Presentation

from worker.celery_app import (
    add_credits_to_slide,
    remove_watermarks_from_masters,
    generate_and_add_speaker_notes,
)


class DummyParent:
    def __init__(self, shapes):
        self._shapes = shapes

    def remove(self, element):
        self._shapes.remove(element.shape)


class DummyElement:
    def __init__(self, parent, shape):
        self._parent = parent
        self.shape = shape

    def getparent(self):
        return self._parent


class DummyShape:
    def __init__(self, text, parent):
        self.has_text_frame = True
        self.text = text
        self.element = DummyElement(parent, self)


class DummyMaster:
    def __init__(self, texts):
        self.shapes = []
        parent = DummyParent(self.shapes)
        for text in texts:
            shape = DummyShape(text, parent)
            self.shapes.append(shape)


class DummyPresentation:
    def __init__(self, master_texts):
        self.slide_masters = [DummyMaster(master_texts)]


@pytest.fixture
def sample_presentation():
    prs = Presentation()
    for _ in range(3):
        prs.slides.add_slide(prs.slide_layouts[6])
    return prs


@pytest.fixture
def ppt_with_watermark():
    return DummyPresentation(["Confidential", "Regular text"])


@pytest.fixture
def slide_with_text():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Q3 Results"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Sales are up 20% year-over-year."
    return slide


def test_add_credits_to_slide_with_custom_text(sample_presentation):
    slide = sample_presentation.slides[0]
    slide_width = sample_presentation.slide_width
    slide_height = sample_presentation.slide_height

    custom_text = "My Custom Brand"
    custom_url = "https://mybrand.com"

    add_credits_to_slide(slide, slide_width, slide_height, custom_text, custom_url)

    assert len(slide.shapes) == 1
    new_shape = slide.shapes[0]
    assert new_shape.text_frame.text == custom_text

    hyperlink = new_shape.text_frame.paragraphs[0].runs[0].hyperlink
    assert hyperlink.address == custom_url


def test_remove_watermarks_from_masters(ppt_with_watermark):
    remove_watermarks_from_masters(ppt_with_watermark)
    shapes = ppt_with_watermark.slide_masters[0].shapes
    assert all("CONFIDENTIAL" not in shape.text.upper() for shape in shapes)


def test_generate_speaker_notes(slide_with_text):
    fake_ai_response_text = "Key takeaway: Strong performance in Q3."

    with patch('worker.celery_app.model.generate_content') as mock_generate_content:
        mock_response = MagicMock()
        mock_response.text = fake_ai_response_text
        mock_generate_content.return_value = mock_response

        generate_and_add_speaker_notes(slide_with_text)

        mock_generate_content.assert_called_once()
        prompt_sent_to_ai = mock_generate_content.call_args[0][0]
        assert "Q3 Results" in prompt_sent_to_ai
        assert "Sales are up 20%" in prompt_sent_to_ai

        notes_text = slide_with_text.notes_slide.notes_text_frame.text
        assert notes_text == fake_ai_response_text
