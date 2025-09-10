import json
import random
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from PIL import Image

# --- Constants ---
SLIDE_W_IN, SLIDE_H_IN = 10.0, 7.5
MARGIN_IN, TITLE_H_IN = 0.5, 0.8
DPI = 96
BULLET_PT = 24
USABLE_W_IN = SLIDE_W_IN - 2 * MARGIN_IN
USABLE_H_IN = SLIDE_H_IN - 2 * MARGIN_IN - TITLE_H_IN
SLIDE_W_PX, SLIDE_H_PX = SLIDE_W_IN * DPI, SLIDE_H_IN * DPI
W_THRESH_PX, H_THRESH_PX = SLIDE_W_PX / 2, SLIDE_H_PX / 2

# --- Helper Functions ---

def classify_image(path: Path):
    """Classifies an image into one of four size categories."""
    if not path:
        return None
    try:
        with Image.open(path) as im:
            w_px, h_px = im.size
        h_cls = "small" if w_px <= W_THRESH_PX else "large"
        v_cls = "small" if h_px <= H_THRESH_PX else "large"
        return f"{h_cls}-{v_cls}"
    except (FileNotFoundError, OSError):
        return None

def add_bullets(slide, bullets, x_in, y_in, w_in, h_in):
    """
    Adds a bulleted (and optionally nested) list to a slide.
    - Leading indentation (spaces or tabs) determines level (2 spaces = 1 level).
    - Leading '-', '*', '•', '–', or numeric '1.' markers are stripped.
    """
    textbox = slide.shapes.add_textbox(Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in))
    tf = textbox.text_frame
    tf.clear()
    tf.word_wrap = True

    # Choose bullet glyph per level (feel free to tweak)
    level_bullets = {0: "•", 1: "◦", 2: "▪", 3: "–", 4: "•"}

    def compute_level_and_text(line: str):
        # Strip trailing spaces
        raw = line.rstrip()
        # Count indentation: tabs = 2 spaces, then groups of 2 spaces => 1 level
        leading = len(raw) - len(raw.lstrip(' \t'))
        # normalize tabs as 2 spaces each
        norm_leading = raw[:leading].replace('\t', '  ')
        spaces = len(norm_leading)
        level = min(spaces // 2, 4)
        text = raw.lstrip(' \t')

        # Strip common list markers
        # e.g., "- item", "* item", "• item", "– item", "1. item", "2) item"
        for marker in ("- ", "* ", "• ", "– "):
            if text.startswith(marker):
                text = text[len(marker):].lstrip()
                break
        else:
            # numeric markers
            import re
            m = re.match(r"^\d+[\.\)]\s+", text)
            if m:
                text = text[m.end():].lstrip()

        return level, text

    for idx, raw in enumerate(bullets):
        level, clean_text = compute_level_and_text(str(raw))

        # ensure at least one paragraph exists
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Manually add a bullet glyph; python-pptx doesn’t toggle bullets for textboxes
        glyph = level_bullets.get(level, "•")
        p.text = f"{glyph} {clean_text}"

        # Formatting
        p.level = level                # controls indent
        p.font.size = Pt(BULLET_PT)    # global const you defined
        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        p.space_after = Pt(4)          # a little breathing room

def add_image_scaled(slide, img_path, x_in, y_in, max_w_in, max_h_in):
    """Adds an image, scaling it to fit the bounds while preserving aspect ratio."""
    with Image.open(img_path) as img:
        img_w_px, img_h_px = img.size
        max_w_px = max_w_in * DPI
        max_h_px = max_h_in * DPI
        
        scale = min(max_w_px / img_w_px, max_h_px / img_h_px)
        
        final_w_in = (img_w_px * scale) / DPI
        final_h_in = (img_h_px * scale) / DPI
        
        x_centered = x_in + (max_w_in - final_w_in) / 2
        y_centered = y_in + (max_h_in - final_h_in) / 2

        slide.shapes.add_picture(
            str(img_path), Inches(x_centered), Inches(y_centered), 
            width=Inches(final_w_in), height=Inches(final_h_in)
        )

# --- Main Builder Function ---

def build_presentation_from_plan(job_dir: Path, output_filename: str):
    json_path = job_dir / "slides.json"
    output_path = job_dir / output_filename

    with open(json_path, 'r', encoding='utf-8') as f:
        specs = json.load(f)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank_layout = prs.slide_layouts[6]

    image_files = sorted([f for f in job_dir.iterdir() if f.suffix.lower() in ['.png', '.jpg', '.jpeg']])

    for i, spec in enumerate(specs):
        slide = prs.slides.add_slide(blank_layout)
        
        title_shape = slide.shapes.add_textbox(Inches(MARGIN_IN), Inches(MARGIN_IN), Inches(USABLE_W_IN), Inches(TITLE_H_IN))
        title_shape.text_frame.text = spec.get("slide_title", " ")
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        
        content_top_in = MARGIN_IN + TITLE_H_IN
        img_path = image_files[i] if i < len(image_files) else None
        classification = classify_image(img_path)
        bullets = spec.get("slide_content", [])
        
        if not classification:
            add_bullets(slide, bullets, MARGIN_IN, content_top_in, USABLE_W_IN, USABLE_H_IN)
        elif classification in ["small-small", "small-large"]:
            variant = random.choice(["left", "right"])
            panel_w = USABLE_W_IN / 2.1
            gap = 0.2
            text_w = USABLE_W_IN - panel_w - gap
            if variant == "left":
                add_image_scaled(slide, img_path, MARGIN_IN, content_top_in, panel_w, USABLE_H_IN)
                add_bullets(slide, bullets, MARGIN_IN + panel_w + gap, content_top_in, text_w, USABLE_H_IN)
            else:
                add_bullets(slide, bullets, MARGIN_IN, content_top_in, text_w, USABLE_H_IN)
                add_image_scaled(slide, img_path, MARGIN_IN + text_w + gap, content_top_in, panel_w, USABLE_H_IN)
        elif classification == "large-small":
            panel_h = USABLE_H_IN / 2.1
            gap = 0.2
            text_h = USABLE_H_IN - panel_h - gap
            add_image_scaled(slide, img_path, MARGIN_IN, content_top_in, USABLE_W_IN, panel_h)
            add_bullets(slide, bullets, MARGIN_IN, content_top_in + panel_h + gap, USABLE_W_IN, text_h)
        elif classification == "large-large":
            # If there are many bullets, use a left-right (side-by-side) layout to avoid overflow.
            if len(bullets) > 4:
                gap = 0.2
                # Give the image a bit more width since it's large in both dims
                panel_w = USABLE_W_IN * 0.55
                text_w = USABLE_W_IN - panel_w - gap

                # Randomize which side the image goes on (optional)
                variant = random.choice(["left", "right"])
                if variant == "left":
                    # Image on left, bullets on right
                    add_image_scaled(slide, img_path, MARGIN_IN, content_top_in, panel_w, USABLE_H_IN)
                    add_bullets(slide, bullets, MARGIN_IN + panel_w + gap, content_top_in, text_w, USABLE_H_IN)
                else:
                    # Bullets on left, image on right
                    add_bullets(slide, bullets, MARGIN_IN, content_top_in, text_w, USABLE_H_IN)
                    add_image_scaled(slide, img_path, MARGIN_IN + text_w + gap, content_top_in, panel_w, USABLE_H_IN)
            else:
                # Fewer bullets: keep original top-bottom layout
                img_h = USABLE_H_IN * 0.7
                text_h = USABLE_H_IN - img_h
                add_image_scaled(slide, img_path, MARGIN_IN, content_top_in, USABLE_W_IN, img_h)
                add_bullets(slide, bullets, MARGIN_IN, content_top_in + img_h, USABLE_W_IN, text_h)

        if spec.get("speaker_notes"):
            slide.notes_slide.notes_text_frame.text = spec.get("speaker_notes")
            
    prs.save(output_path)
    return output_path