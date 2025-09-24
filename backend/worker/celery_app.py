import os
import io
import json
import base64
import shutil
from pathlib import Path
from typing import Any, Iterable
from celery import Celery
from pptx import Presentation
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageDraw
import imagehash
import requests
from google.cloud import storage
from google.auth.exceptions import DefaultCredentialsError

from config import settings
from config.storage import LocalStorageClient


class _NoopModel:
    def generate_content(self, *_args, **_kwargs):
        raise RuntimeError("Speaker notes model is not configured")


def _create_storage_client() -> storage.Client:
    """Instantiate a storage client with graceful fallback for local/dev runs."""
    if settings.use_local_storage:
        return LocalStorageClient()
    try:
        return storage.Client()
    except DefaultCredentialsError:
        if settings.is_development:
            return LocalStorageClient()
        raise


# Import other project modules
from .creator_logic import (
    extract_text_from_document,
    generate_content_for_batch,
    generate_outline_from_text,
    genai_client,
    model as text_generation_model,
)
from .ppt_builder import build_presentation_from_plan

# --- Configuration ---
GCS_BUCKET_NAME = settings.gcs_bucket_name
storage_client = _create_storage_client()
LOGO_PATH = "temp/logo.png"  # Default logo path if none is provided
WATERMARK_KEYWORDS = ["CONFIDENTIAL", "DRAFT", "INTERNAL USE"]

speaker_notes_model = text_generation_model
# Backwards compatibility for modules/tests referencing `worker.celery_app.model`
model = speaker_notes_model

# --- Initialize Celery ---
celery_backend = settings.celery_backend_url
celery = Celery("tasks", broker=settings.celery_broker_url, backend=celery_backend or None)
celery.conf.update(
    task_serializer="json",
    accept_content=["json"],
    result_serializer="json",
    result_expires=settings.celery_result_expires,
)
if not celery_backend:
    celery.conf.update(result_backend=None, task_ignore_result=True)

# --- GCS Helper Functions ---
def download_blob(blob_name, destination_file_name):
    bucket = storage_client.bucket(GCS_BUCKET_NAME)
    blob = bucket.blob(blob_name)
    blob.download_to_filename(destination_file_name)

def upload_blob(source_file_name, destination_blob_name):
    bucket = storage_client.bucket(GCS_BUCKET_NAME)
    blob = bucket.blob(destination_blob_name)
    blob.upload_from_filename(source_file_name)

def list_blobs(prefix):
    bucket = storage_client.bucket(GCS_BUCKET_NAME)
    return bucket.list_blobs(prefix=prefix)


def _metadata_key(job_id: str) -> str:
    return f"{job_id}/metadata.json"


def _load_job_metadata(job_id: str) -> dict:
    bucket = storage_client.bucket(GCS_BUCKET_NAME)
    blob = bucket.blob(_metadata_key(job_id))
    if not blob.exists():
        return {}
    try:
        return json.loads(blob.download_as_text())
    except Exception:
        return {}


def _store_job_metadata(job_id: str, metadata: dict) -> None:
    bucket = storage_client.bucket(GCS_BUCKET_NAME)
    bucket.blob(_metadata_key(job_id)).upload_from_string(
        json.dumps(metadata, indent=2), content_type="application/json"
    )


def _extension_from_mime(mime: str | None) -> str:
    if not mime or "/" not in mime:
        return "png"
    return mime.split("/")[-1] or "png"


def _auto_generate_images(job_id: str, strategy: str, outline: list[dict[str, Any]], base_dir: Path, client: Any | None) -> list[dict[str, Any]]:
    if not outline:
        raise RuntimeError("Outline not available for image generation")

    records: list[dict[str, Any]] = []
    bucket = storage_client.bucket(GCS_BUCKET_NAME)

    for index, slide in enumerate(outline, start=1):
        title = slide.get("slide_title") or f"Slide {index}"
        if strategy == "unsplash":
            image_bytes, extension = _fetch_unsplash_image(slide)
        else:
            image_bytes, extension = _generate_line_art_image(slide, client)

        filename = _safe_filename(title, index, extension)
        local_path = base_dir / filename
        with open(local_path, "wb") as fh:
            fh.write(image_bytes)

        blob_path = f"auto-images/{filename}"
        bucket.blob(f"{job_id}/{blob_path}").upload_from_filename(str(local_path))
        records.append({
            "blob_path": blob_path,
            "filename": filename,
            "local_path": str(local_path),
        })

    return records


def _outline_blob_name(job_id: str) -> str:
    return f"{job_id}/outline.json"


def _load_outline(job_id: str) -> list[dict[str, Any]]:
    bucket = storage_client.bucket(GCS_BUCKET_NAME)
    blob = bucket.blob(_outline_blob_name(job_id))
    if not blob.exists():
        return []
    try:
        return json.loads(blob.download_as_text())
    except Exception:
        return []


def _placeholder_image_bytes(title: str) -> bytes:
    image = Image.new("RGB", (1280, 720), color="white")
    drawer = ImageDraw.Draw(image)
    text = f"{title[:60] or 'Slide'}"
    drawer.text((40, 340), text, fill="black")
    buffer = io.BytesIO()
    image.save(buffer, format="PNG")
    return buffer.getvalue()


def _safe_filename(base: str, index: int, extension: str) -> str:
    slug = ''.join(ch.lower() if ch.isalnum() else '-' for ch in base)[:50]
    slug = slug.strip('-') or f"slide-{index}"
    slug = '-'.join(filter(None, slug.split('-')))
    return f"{slug or f'slide-{index}'}-{index}.{extension}"


def _fetch_unsplash_image(slide: dict[str, Any]) -> tuple[bytes, str]:
    query_parts: list[str] = []
    if title := slide.get("slide_title"):
        query_parts.append(title)
    bullets: Iterable[str] | None = slide.get("bullet_outline")
    if bullets:
        first_bullet = next(iter(bullets), None)
        if first_bullet:
            query_parts.append(first_bullet)
    query = ' '.join(query_parts) or "presentation"

    if not settings.unsplash_access_key:
        print("Unsplash access key not configured; using placeholder image")
        return _placeholder_image_bytes(slide.get("slide_title", "Slide")), "png"

    headers = {"Authorization": f"Client-ID {settings.unsplash_access_key}"}
    params = {"query": query, "orientation": settings.unsplash_orientation}
    try:
        response = requests.get(
            "https://api.unsplash.com/photos/random",
            headers=headers,
            params=params,
            timeout=10,
        )
        response.raise_for_status()
        payload = response.json()
        image_url = (
            payload.get("urls", {}).get("regular")
            or payload.get("urls", {}).get("full")
            or payload.get("urls", {}).get("small")
        )
        if not image_url:
            raise RuntimeError("Unsplash response missing image URL")

        image_response = requests.get(image_url, timeout=15)
        image_response.raise_for_status()
        return image_response.content, "jpg"
    except Exception as exc:
        print(f"Unsplash image fetch failed: {exc}")
        return _placeholder_image_bytes(slide.get("slide_title", "Slide")), "png"


def _generate_line_art_image(slide: dict[str, Any], client: Any | None) -> tuple[bytes, str]:
    prompt = (
        "Create a clean black-and-white line art illustration suitable for a presentation slide. "
        f"Focus on the concept: {slide.get('slide_title', 'concept')}. "
        "Use minimal shading, crisp outlines, and no color."
    )
    if bullets := slide.get("bullet_outline"):
        prompt += " Key points: " + ", ".join(bullets[:2]) + "."

    if client is None:
        return _placeholder_image_bytes(slide.get("slide_title", "Slide")), "png"

    try:
        response = client.models.generate_images(
            model=settings.gemini_image_model,
            prompt=prompt,
            config={
                "aspectRatio": "16:9",
                "numberOfImages": 1,
                "outputMimeType": "image/png",
            },
        )
        generated = getattr(response, "generated_images", None) or []
        for generated_image in generated:
            image = getattr(generated_image, "image", None)
            if not image:
                continue
            # Prefer direct bytes if available
            data = getattr(image, "image_bytes", None)
            if data is None:
                data = getattr(image, "imageBytes", None)
            if data is None:
                # Some SDKs use gcsUri instead
                gcs_uri = getattr(image, "gcs_uri", None) or getattr(image, "gcsUri", None)
                if gcs_uri:
                    try:
                        response = requests.get(gcs_uri, timeout=30)
                        if response.ok:
                            return response.content, _extension_from_mime(response.headers.get("Content-Type"))
                    except Exception as exc:
                        print(f"Failed to download Gemini image from GCS URI {gcs_uri}: {exc}")
                continue

            if isinstance(data, str):
                image_bytes = base64.b64decode(data)
            else:
                image_bytes = data

            mime = getattr(image, "mime_type", None) or getattr(image, "mimeType", None) or "image/png"
            return image_bytes, _extension_from_mime(mime)
    except Exception as exc:
        print(f"Gemini image generation failed: {exc}")

    return _placeholder_image_bytes(slide.get("slide_title", "Slide")), "png"

# --- Other Business Logic (Full versions) ---
def chunks(lst, n):
    for i in range(0, len(lst), n):
        yield lst[i:i + n]

def _iter_picture_shapes(container):
    if not hasattr(container, "shapes"): return
    for shape in list(container.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            yield shape
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for inner in _iter_picture_shapes(shape):
                yield inner

def _image_hash_for_shape(shape):
    try:
        with Image.open(io.BytesIO(shape.image.blob)).convert("RGB") as im:
            return imagehash.phash(im)
    except Exception:
        return None

def remove_frequent_images(prs: Presentation, min_occurrences: int, hash_tolerance: int):
    all_pics_and_hashes = []
    for master in prs.slide_masters:
        for pic in _iter_picture_shapes(master):
            if h := _image_hash_for_shape(pic):
                all_pics_and_hashes.append({'shape': pic, 'hash': h})
    for layout in prs.slide_layouts:
        for pic in _iter_picture_shapes(layout):
            if h := _image_hash_for_shape(pic):
                all_pics_and_hashes.append({'shape': pic, 'hash': h})
    for slide in prs.slides:
        for pic in _iter_picture_shapes(slide):
            if h := _image_hash_for_shape(pic):
                all_pics_and_hashes.append({'shape': pic, 'hash': h})
    if not all_pics_and_hashes: return
    hash_clusters = {}
    for item in all_pics_and_hashes:
        found_cluster = False
        for h_key in hash_clusters:
            if item['hash'] - h_key <= hash_tolerance:
                hash_clusters[h_key].append(item['shape'])
                found_cluster = True
                break
        if not found_cluster:
            hash_clusters[item['hash']] = [item['shape']]
    shapes_to_delete = []
    for h, shapes in hash_clusters.items():
        if len(shapes) >= min_occurrences:
            shapes_to_delete.extend(shapes)
    for shape in shapes_to_delete:
        sp = shape.element
        sp.getparent().remove(sp)

def extract_text_from_slide(slide: Slide) -> str:
    slide_texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                slide_texts.append(text)
    return "\n".join(slide_texts)

def generate_and_add_speaker_notes(slide: Slide):
    try:
        slide_text = extract_text_from_slide(slide)
        if not slide_text: return
        if isinstance(speaker_notes_model, _NoopModel):
            raise RuntimeError("Speaker notes model is not configured")
        prompt = f"Generate a concise, professional speaker note for a presentation slide with the following content:\n\n---\n{slide_text}\n---"
        response = speaker_notes_model.generate_content(prompt)
        if response.text:
            slide.notes_slide.notes_text_frame.text = response.text
    except Exception as e:
        slide.notes_slide.notes_text_frame.text = f"Could not generate speaker notes: {e}"

def remove_watermarks_from_masters(prs: Presentation):
    for master in prs.slide_masters:
        shapes_to_delete = [
            shape for shape in master.shapes
            if shape.has_text_frame and any(keyword in shape.text.upper() for keyword in WATERMARK_KEYWORDS)
        ]
        for shape in shapes_to_delete:
            sp = shape.element
            sp.getparent().remove(sp)

def add_logo(slide: Slide, logo_path: str):
    if logo_path and os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(0.2), Inches(0.2), width=Inches(1.0))

def add_credits_to_slide(slide: Slide, slide_width, slide_height, text: str, url: str):
    textbox = slide.shapes.add_textbox(
        left=slide_width - Inches(2.6), top=slide_height - Inches(0.5),
        width=Inches(2.5), height=Inches(0.4)
    )
    tf = textbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT
    run = p.add_run()
    run.text = text
    run.hyperlink.address = url
    font = run.font
    font.size = Pt(10)
    font.color.rgb = RGBColor(150, 150, 150)

# --- Celery Tasks ---
@celery.task(name="enhance_ppt_task")
def enhance_ppt_task(input_blob: str, output_blob: str, logo_blob: str = None, credits_text: str = None):
    job_id = Path(input_blob).parts[0]
    local_job_dir = Path("/tmp") / job_id
    local_job_dir.mkdir(parents=True, exist_ok=True)
    try:
        local_input_path = local_job_dir / Path(input_blob).name
        download_blob(input_blob, str(local_input_path))
        
        local_logo_path = None
        if logo_blob:
            local_logo_path = local_job_dir / Path(logo_blob).name
            download_blob(logo_blob, str(local_logo_path))
        
        prs = Presentation(local_input_path)
        final_credits_text = credits_text if credits_text else "Processed by PPT Studio"
        final_credits_url = "https://mybrand.com" if credits_text else "https://www.example.com"
        final_logo_path = str(local_logo_path) if local_logo_path and local_logo_path.exists() else LOGO_PATH
        
        remove_watermarks_from_masters(prs)
        remove_frequent_images(prs, min_occurrences=3, hash_tolerance=5)

        for slide in prs.slides:
            shapes_to_delete_text = [
                shape for shape in slide.shapes
                if shape.has_text_frame and any(keyword in shape.text.upper() for keyword in WATERMARK_KEYWORDS)
            ]
            for shape in shapes_to_delete_text:
                sp = shape.element
                sp.getparent().remove(sp)
            add_logo(slide, final_logo_path)
            add_credits_to_slide(slide, prs.slide_width, prs.slide_height, final_credits_text, final_credits_url)
            generate_and_add_speaker_notes(slide)
        
        local_output_path = local_job_dir / Path(output_blob).name
        prs.save(str(local_output_path))
        upload_blob(str(local_output_path), output_blob)
        return {"status": "complete", "output_blob": output_blob}
    finally:
        shutil.rmtree(local_job_dir, ignore_errors=True)

@celery.task(name="generate_outline_task")
def generate_outline_task(job_id: str, source_blob: str, slide_count: int | None = None, **_extra):
    local_job_dir = Path("/tmp") / job_id / "outline"
    local_job_dir.mkdir(parents=True, exist_ok=True)
    try:
        metadata = _load_job_metadata(job_id)
        if slide_count is None:
            raw_count = metadata.get("desired_slide_count") if metadata else None
            try:
                slide_count = max(1, int(raw_count)) if raw_count is not None else None
            except (TypeError, ValueError):
                slide_count = None

        local_source_path = local_job_dir / Path(source_blob).name
        download_blob(f"{job_id}/{source_blob}", str(local_source_path))
        source_text = extract_text_from_document(str(local_source_path))
        outline = generate_outline_from_text(source_text, desired_slide_count=slide_count)

        if slide_count is not None:
            try:
                target = max(1, int(slide_count))
            except (TypeError, ValueError):
                target = None
            if target:
                if len(outline) > target:
                    outline = outline[:target]
                elif len(outline) < target:
                    while len(outline) < target:
                        outline.append({
                            "slide_title": f"Slide {len(outline) + 1}",
                            "bullet_outline": ["Outline content to be refined."],
                        })
                slide_count = target

        outline_path = local_job_dir / "outline.json"
        with open(outline_path, "w") as fh:
            json.dump(outline, fh, indent=2)
        upload_blob(str(outline_path), f"{job_id}/outline.json")

        metadata = metadata or {}
        if "source_blob" not in metadata:
            metadata["source_blob"] = source_blob
        metadata.update({
            "outline_blob": "outline.json",
            "outline_generated": True,
        })
        if slide_count:
            metadata["desired_slide_count"] = slide_count
        _store_job_metadata(job_id, metadata)

        return {
            "status": "complete",
            "outline": outline,
            "desired_slide_count": slide_count,
        }
    finally:
        shutil.rmtree(local_job_dir, ignore_errors=True)


@celery.task(name="generate_slide_plan_task")
def generate_slide_plan_task(
    job_id: str,
    image_records: list | None = None,
    source_blob: str | None = None,
    image_strategy: str | None = None,
    **_ignore_extra_kwargs,
):
    local_job_dir = Path("/tmp") / job_id
    local_job_dir.mkdir(parents=True, exist_ok=True)
    try:
        image_records = image_records or []
        metadata = _load_job_metadata(job_id)
        if not image_strategy:
            image_strategy = metadata.get("image_strategy", "uploaded")
        if not source_blob:
            source_blob = metadata.get("source_blob")

        image_blob_paths = []
        display_names = []
        local_image_paths = []

        if image_strategy in {"unsplash", "gemini_line_art"}:
            outline = _load_outline(job_id)
            try:
                auto_records = _auto_generate_images(
                    job_id,
                    image_strategy,
                    outline,
                    local_job_dir,
                    genai_client,
                )
                image_records.extend(auto_records)
                metadata["auto_image_strategy"] = image_strategy
                metadata["auto_images"] = [rec["blob_path"] for rec in auto_records]
            except Exception as exc:
                return {"error": f"Failed to create images automatically: {exc}"}

        for record in image_records:
            if isinstance(record, dict):
                blob_rel = record.get("blob_path") or record.get("blob") or record.get("path")
                display_name = record.get("filename") or (Path(blob_rel).name if blob_rel else None)
                local_override = record.get("local_path")
            else:
                blob_rel = record
                display_name = Path(blob_rel).name
                local_override = None

            if not blob_rel:
                continue

            display_name = display_name or Path(blob_rel).name
            image_blob_paths.append(blob_rel)
            display_names.append(display_name)

            if local_override and Path(local_override).exists():
                local_image_paths.append(Path(local_override))
            else:
                local_path = local_job_dir / display_name
                download_blob(f"{job_id}/{blob_rel}", str(local_path))
                local_image_paths.append(local_path)

        if not local_image_paths:
            return {"error": "No images were provided for slide planning."}

        if not source_blob:
            job_prefix = f"{job_id}/"
            for blob in list_blobs(job_id):
                rel_path = blob.name[len(job_prefix):] if blob.name.startswith(job_prefix) else blob.name
                rel_name = Path(rel_path).name
                if rel_path in image_blob_paths or rel_name in display_names:
                    continue
                if rel_name in {"outline.json", "slides.json"}:
                    continue
                source_blob = rel_path
                break

        if not source_blob:
            return {"error": "No source document found in GCS."}

        local_source_path = local_job_dir / Path(source_blob).name
        download_blob(f"{job_id}/{source_blob}", str(local_source_path))
        source_text = extract_text_from_document(str(local_source_path))

        slide_plan = generate_content_for_batch(source_text, local_image_paths)
        if not slide_plan:
            return {"error": "Failed to generate a slide plan."}

        for slide, display_name in zip(slide_plan, display_names):
            slide.setdefault("image_filename", display_name)

        plan_path = local_job_dir / "slides.json"
        with open(plan_path, "w") as f:
            json.dump(slide_plan, f, indent=2)
        upload_blob(str(plan_path), f"{job_id}/slides.json")

        metadata.setdefault("source_blob", source_blob)
        metadata.update({
            "last_plan_blob": "slides.json",
            "image_filenames": display_names,
            "image_strategy": image_strategy,
        })
        _store_job_metadata(job_id, metadata)

        return {
            "status": "complete",
            "slide_plan": slide_plan,
            "image_strategy": image_strategy,
        }
    finally:
        shutil.rmtree(local_job_dir, ignore_errors=True)

@celery.task(name="build_ppt_from_plan_task")
def build_ppt_from_plan_task(job_id: str):
    local_job_dir = Path("/tmp") / job_id
    local_job_dir.mkdir(parents=True, exist_ok=True)
    try:
        blobs = list_blobs(job_id)
        for blob in blobs:
            download_blob(blob.name, str(local_job_dir / Path(blob.name).name))
        
        output_filename = "presentation.pptx"
        local_output_path = build_presentation_from_plan(local_job_dir, output_filename)
        upload_blob(str(local_output_path), f"{job_id}/{output_filename}")
        return {"status": "complete", "output_file": f"{job_id}/{output_filename}"}
    finally:
        shutil.rmtree(local_job_dir, ignore_errors=True)
