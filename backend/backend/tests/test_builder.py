import json # <-- Add this import
from pathlib import Path
from pptx import Presentation
from PIL import Image
from backend.worker.ppt_builder import build_presentation_from_plan

def test_build_presentation_from_plan(tmp_path):
    """
    Tests that a .pptx file is correctly generated from a job directory
    containing a slides.json and an image.
    """
    # 1. Setup: Create a fake job directory with test assets
    job_dir = tmp_path
    
    # Create a fake slides.json
    slide_plan = [{
        "slide_title": "Test Slide",
        "slide_content": ["- Bullet 1", "- Bullet 2"],
        "speaker_notes": "This is a test note."
    }]
    with open(job_dir / "slides.json", "w") as f:
        json.dump(slide_plan, f)
        
    # Create a fake image file
    fake_image = Image.new('RGB', (100, 100), color = 'red')
    fake_image_path = job_dir / "test_image.png"
    fake_image.save(fake_image_path)
    
    # 2. Execution: Call the function we are testing
    output_filename = "final_presentation.pptx"
    output_path = build_presentation_from_plan(job_dir, output_filename)
    
    # 3. Assertions: Check the output
    assert output_path.exists()
    
    # Optional: Check the content of the generated PPTX
    prs = Presentation(output_path)
    assert len(prs.slides) == 1
    assert prs.slides[0].notes_slide.notes_text_frame.text == "This is a test note."

def test_build_presentation_with_wide_image(tmp_path):
    """
    Tests that a wide (landscape) image results in a top/bottom panel layout.
    """
    # 1. Setup
    job_dir = tmp_path
    slide_plan = [{"slide_title": "Wide Image Test", "slide_content": ["- Landscape image"]}]
    with open(job_dir / "slides.json", "w") as f:
        json.dump(slide_plan, f)
        
    wide_image = Image.new('RGB', (800, 200), color='blue')
    wide_image.save(job_dir / "wide_image.png")
    
    # 2. Execution
    output_path = build_presentation_from_plan(job_dir, "wide_test.pptx")
    
    # 3. Assertions
    prs = Presentation(output_path)
    assert len(prs.slides) == 1
    slide = prs.slides[0]
    assert len(slide.shapes) >= 3 # Title, Image, Textbox
    
    image_shape = None
    for shape in slide.shapes:
        if hasattr(shape, 'image'):
            image_shape = shape
            break
            
    assert image_shape is not None, "Image shape not found on slide"
    assert image_shape.width > image_shape.height * 2